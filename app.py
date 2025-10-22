import io, os, re, json, zipfile
from collections import Counter, defaultdict

import streamlit as st
import fitz  # PyMuPDF
from pptx import Presentation
from PIL import Image
import numpy as np

from sentence_transformers import SentenceTransformer
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

# ---------- UI ----------
st.set_page_config(page_title="ClassMind — AI Study Assistant", page_icon="🧠", layout="wide")
st.title("ClassMind — AI Study Assistant")
st.write("Upload a PDF or PPTX. Get a deck summary (no hallucinations), semantic search, a Mermaid mind map, and Anki flashcards.")

# ---------- Optional OCR (safe fallback) ----------
try:
    import pytesseract
    OCR_AVAILABLE = True
except Exception:
    OCR_AVAILABLE = False

def _clean_text(s: str) -> str:
    if not s: return ""
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()

def _ocr_pil(img: Image.Image) -> str:
    if not OCR_AVAILABLE: return ""
    OCR_LANG = "eng+deu"
    OCR_PSM = "6"
    OCR_OEM = "3"
    cfg = f"--oem {OCR_OEM} --psm {OCR_PSM}"
    try:
        return _clean_text(pytesseract.image_to_string(img, lang=OCR_LANG, config=cfg))
    except Exception:
        return ""

def _pdf_page_to_image(page, dpi=300) -> Image.Image:
    pix = page.get_pixmap(dpi=dpi)
    return Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")

# ---------- Ingestors ----------
def extract_pdf_text(file_bytes, ocr=True, min_chars_no_ocr=25, dpi=250):
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    slides, ocr_count = [], 0
    for i, page in enumerate(doc):
        raw = (page.get_text() or "").strip()
        used_ocr = False
        image_only = False
        if (not raw or len(raw) < min_chars_no_ocr) and ocr and OCR_AVAILABLE:
            img = _pdf_page_to_image(page, dpi=dpi)
            raw = _ocr_pil(img)
            used_ocr = True
            image_only = True
            ocr_count += 1
        slides.append({
            "index": i,
            "text": _clean_text(raw),
            "source": "pdf",
            "ocr_used": used_ocr,
            "image_only": image_only
        })
    return slides, ocr_count, len(slides)

def extract_pptx_text(file_bytes, ocr=True):
    prs = Presentation(io.BytesIO(file_bytes))
    slides, ocr_count = [], 0
    for i, s in enumerate(prs.slides):
        texts, images_ocr = [], []
        has_picture = False
        for shape in s.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
            if getattr(shape, "shape_type", None) == 13:
                has_picture = True
                if ocr and OCR_AVAILABLE:
                    try:
                        img = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                        images_ocr.append(_ocr_pil(img))
                    except Exception:
                        pass
        raw = _clean_text("\n".join(t for t in texts if t))
        used_ocr = False
        image_only = False
        if ocr and (not raw or len(raw) < 25) and has_picture and images_ocr:
            ocr_txt = _clean_text("\n".join(t for t in images_ocr if t))
            if ocr_txt:
                raw = (raw + "\n" + ocr_txt).strip() if raw else ocr_txt
                used_ocr = True
                image_only = (len(_clean_text("\n".join(texts))) < 5)
                ocr_count += 1
        slides.append({
            "index": i,
            "text": raw,
            "source": "pptx",
            "ocr_used": used_ocr,
            "image_only": image_only
        })
    return slides, ocr_count, len(slides)

# ---------- Models (cache) ----------
@st.cache_resource(show_spinner=False)
def load_embedder():
    return SentenceTransformer("sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2")

embedder = load_embedder()

# ---------- Extractive Summary ----------
def extractive_summary_from_slides(slides, max_words=260):
    pages = [s.get("text","") for s in slides if s.get("text")]
    doc = "\n".join(pages)
    sentences = [t.strip() for t in re.split(r'(?<=[.!?])\s+', doc) if len(t.strip().split())>=6]

    def keep_sentence(s):
        if sum(c.isalpha() for c in s)/max(1,len(s)) < 0.5: return False
        if re.search(r'(Chair of|Page \d+|\*{3}Not for sharing\*{3}|https?://|Wikimedia|CC BY|Pixabay)', s, re.I): return False
        return True
    sentences = [s for s in sentences if keep_sentence(s)]
    if not sentences:
        return ""

    stop = """der die das ein eine und oder für von mit ohne zu auf an im in aus ist sind war waren sein auch sowie
the a an and or for from with without into in on of to is are was were be been being this that these those it its by as at if then else when while not
data using use based results method methods model models paper figure table slide page""".split()
    vec = TfidfVectorizer(stop_words=stop, max_df=0.9, min_df=2, ngram_range=(1,2))
    X = vec.fit_transform(sentences)
    scores = X.sum(axis=1).A.ravel()

    sim = cosine_similarity(X, X)
    chosen_idx = []
    remaining = set(range(len(sentences)))
    i0 = int(scores.argmax()); chosen_idx.append(i0); remaining.remove(i0)
    total = len(sentences[i0].split())
    while remaining and total < max_words:
        best_i, best_v = None, -1e9
        for i in remaining:
            rel = scores[i]
            div = 0 if not chosen_idx else max(sim[i, chosen_idx])
            val = 0.7 * rel - 0.3 * div
            if val > best_v and total + len(sentences[i].split()) <= max_words + 40:
                best_v, best_i = val, i
        if best_i is None: break
        chosen_idx.append(best_i); remaining.remove(best_i)
        total += len(sentences[best_i].split())

    out = " ".join(sentences[i] for i in sorted(chosen_idx))
    return re.sub(r"\s+", " ", out).strip()

# ---------- Mind map ----------
def top_keywords(text, k=3):
    words = re.findall(r"[A-Za-zÄÖÜäöüß][A-Za-zÄÖÜäöüß\-]{2,}", text or "")
    stop = set("""der die das ein eine und oder für von mit ohne zu auf an im in aus ist sind war waren sein auch sowie
the a an and or for from with without into in on of to is are was were be been being this that these those it its by as at if then else when while not""".split())
    words = [w.lower() for w in words if w.lower() not in stop]
    return [w for w,_ in Counter(words).most_common(k)]

def build_mindmap(slide_summaries, title="Deck"):
    kw_to_slides = defaultdict(list)
    for s in slide_summaries:
        base = s.get("summary") or s.get("text","")
        for kw in top_keywords(base, 3):
            kw_to_slides[kw].append(s["index"])
    top_kws = list(kw_to_slides.keys())[:12]
    lines = ["mindmap", f"  root(({title}))"]
    for kw in top_kws:
        lines.append(f"    {kw}({kw})")
        for idx in sorted(kw_to_slides[kw])[:6]:
            lines.append(f"      s{idx}([Slide {idx}])")
    return "\n".join(lines)

# ---------- Flashcards ----------
def split_points(text):
    bullets = re.findall(r"(?:^|\n)[\-\•\*]\s*(.+)", text or "")
    if bullets: return [b.strip() for b in bullets if len(b.strip())>3][:4]
    sents = re.split(r"(?<=[.!?])\s+", text or "")
    return [s for s in sents if len(s.split())>=6][:3]

def detect_de(text):
    t=(text or "").lower()
    if any(c in t for c in "äöüÄÖÜß"): return True
    de_hits=sum(w in t for w in ["und","der","die","das","ist","nicht","ein","eine","oder","auch"])
    en_hits=sum(w in t for w in ["and","the","is","are","not","or","also","a","an"])
    return de_hits>=en_hits

def make_flashcards(slide_summaries):
    rows=[]
    for s in slide_summaries:
        base=s.get("summary") or s.get("text","")
        if not base: continue
        de=detect_de(base)
        items=split_points(base)
        if items:
            for j,it in enumerate(items,1):
                q = f"Nenne einen Kerngedanken (Folie {s['index']}, Punkt {j}):" if de else f"Name one key idea (slide {s['index']}, point {j}):"
                rows.append([q, it, f"ClassMind;{'DE' if de else 'EN'}"])
        else:
            q = f"Was ist die Kernaussage von Folie {s['index']}?" if de else f"What is the main idea of slide {s['index']}?"
            rows.append([q, base, f"ClassMind;{'DE' if de else 'EN'}"])
    return rows

# ---------- File uploader ----------
uploaded = st.file_uploader("Upload a PDF or PPTX", type=["pdf","pptx"])
if uploaded:
    ext = uploaded.name.lower().split(".")[-1]
    file_bytes = uploaded.read()

    with st.spinner("Extracting slides..."):
        if ext == "pdf":
            slides, ocrn, nslides = extract_pdf_text(file_bytes, ocr=True)
        else:
            slides, ocrn, nslides = extract_pptx_text(file_bytes, ocr=True)
    st.success(f"Slides: {nslides} • OCR used on {ocrn} slide(s) • OCR available: {OCR_AVAILABLE}")

    # Per-slide "summaries": we reuse the raw text to keep it fast; optional: add a short heuristic
    slide_summaries = []
    for s in slides:
        txt = s["text"]
        # tiny heuristic: shorten long slides
        short = txt
        if txt and len(txt.split()) > 80:
            # keep first ~2 sentences as 'summary'
            parts = re.split(r"(?<=[.!?])\s+", txt)
            short = " ".join(parts[:2]).strip()
        slide_summaries.append({**s, "summary": short})

    # Extractive deck summary
    with st.spinner("Summarizing deck (extractive, no hallucinations)..."):
        deck_summary = extractive_summary_from_slides(slides)
    st.subheader("Deck Summary")
    st.write(deck_summary if deck_summary else "_No text detected._")

    # Embeddings + search
    st.subheader("Semantic Search")
    corpus = [f"Slide {s['index']}: {s['summary'] or s['text']}" for s in slide_summaries]
    embs = embedder.encode(corpus, convert_to_numpy=True, normalize_embeddings=True)
    query = st.text_input("Search the deck (DE/EN)", "")
    if query:
        q = embedder.encode([query], convert_to_numpy=True, normalize_embeddings=True)
        sims = (embs @ q.T).ravel()
        top = sims.argsort()[::-1][:5]
        for i in top:
            st.markdown(f"**Slide {slide_summaries[i]['index']}** — score {float(sims[i]):.3f}")
            st.write((slide_summaries[i]["summary"] or slide_summaries[i]["text"])[:400])

    # Mind map
    st.subheader("Mind Map (Mermaid)")
    mmd = build_mindmap(slide_summaries, title=uploaded.name)
    st.code(mmd, language="markdown")
    st.download_button("Download mindmap.mmd", mmd, file_name="mindmap.mmd")

    # Flashcards
    st.subheader("Flashcards (Anki CSV)")
    rows = make_flashcards(slide_summaries)
    import csv
    import io as _io
    csv_buf = _io.StringIO()
    w = csv.writer(csv_buf); w.writerow(["Question","Answer","Tags"]); w.writerows(rows)
    st.download_button("Download flashcards.csv", csv_buf.getvalue(), file_name="flashcards.csv")

    # Bundle all outputs as a zip
    zbuf = _io.BytesIO()
    with zipfile.ZipFile(zbuf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("deck_summary.txt", deck_summary or "")
        zf.writestr("mindmap.mmd", mmd)
        zf.writestr("flashcards.csv", csv_buf.getvalue())
        zf.writestr("slide_summaries.json", json.dumps(slide_summaries, ensure_ascii=False, indent=2))
    st.download_button("Download all (zip)", zbuf.getvalue(), file_name="classmind_outputs.zip")

else:
    st.info("Upload a file to begin. Supported: PDF, PPTX.")

# Footer
with st.sidebar:
    st.markdown("### About")
    st.write("ClassMind turns lecture slides into study assets — **no paid APIs**. This demo uses extractive summarization for speed and reliability.")
    st.write("Tip: Paste `mindmap.mmd` into https://mermaid.live to render a PNG/SVG.")
